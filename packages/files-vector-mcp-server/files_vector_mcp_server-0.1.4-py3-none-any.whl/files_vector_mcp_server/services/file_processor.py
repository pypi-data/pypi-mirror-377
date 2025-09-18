import os
import time
import logging
import hashlib
import asyncio
from datetime import datetime
from typing import List, Dict, Optional, Any

logger = logging.getLogger(__name__)

class FileProcessor:
    """文件处理服务类，负责文件内容读取、分块和向量化"""
    
    def __init__(self, db_service, openai_client, config):
        """初始化文件处理器"""
        self.db = db_service
        self.openai_client = openai_client
        self.config = config
        self.topic_directory_map = self._build_topic_map()
        
        # 检查MinerU API配置，两个参数都配置才启用
        self.use_mineru = bool(config.mineru_api_key and config.mineru_api_url)
        
        if self.use_mineru:
            logger.info(f"已启用MinerU API支持，服务地址: {config.mineru_api_url}")
        else:
            logger.warning("MinerU API未完全配置，将禁用相关功能")

    def _build_topic_map(self) -> Dict[str, str]:
        """构建目录到主题的映射"""
        topic_map = {}
        for topic, directories in self.config.watch_topics.items():
            for directory in directories:
                abs_dir = os.path.abspath(directory)
                topic_map[abs_dir] = topic
        return topic_map

    def _get_file_topic(self, file_path: str) -> str:
        """确定文件所属主题"""
        abs_file_path = os.path.abspath(file_path)
        for directory, topic in self.topic_directory_map.items():
            if abs_file_path.startswith(directory):
                return topic
        logger.warning(f"无法确定文件 {file_path} 的主题，将使用'未分类'主题")
        return "未分类"

    def _get_file_hash(self, file_path: str) -> str:
        """计算文件内容的MD5哈希"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def _read_file_content(self, file_path: str) -> str:
        """读取不同类型文件的文本内容"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            content = ""
            
            # Markdown文件
            if ext == ".md":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            
            # PowerPoint文件
            elif ext in [".ppt", ".pptx"]:
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    content_parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                content_parts.append(shape.text.strip())
                    content = "\n\n".join(content_parts)
                except ImportError:
                    logger.error("python-pptx未安装，无法处理PPT/PPTX文件")
                    raise Exception("需要安装python-pptx来处理PPT/PPTX文件: pip install python-pptx")
            
            # 图片文件 - 仅在MinerU配置时处理
            elif ext in [".jpg", ".jpeg", ".png"]:
                if self.use_mineru:
                    content = self._mineru_api_process(file_path, "image")
                else:
                    logger.warning(f"MinerU未配置，跳过图片文件处理: {file_path}")
                    return ""  # 返回空内容，避免后续处理
            
            # PDF文件
            elif ext == ".pdf":
                if self.use_mineru:
                    content = self._mineru_api_process(file_path, "pdf")
                else:
                    # 回退到本地PDF处理
                    try:
                        from PyPDF2 import PdfReader
                        reader = PdfReader(file_path)
                        content = "\n".join([page.extract_text() for page in reader.pages])
                    except ImportError:
                        logger.error("PyPDF2未安装，无法处理PDF文件")
                        raise Exception("需要安装PyPDF2来处理PDF文件: pip install PyPDF2")
            
            # Word文件
            elif ext in [".docx"]:
                try:
                    from docx import Document
                    doc = Document(file_path)
                    content = "\n".join([para.text for para in doc.paragraphs])
                except ImportError:
                    logger.error("python-docx未安装，无法处理DOCX文件")
                    raise Exception("需要安装python-docx来处理DOCX文件: pip install python-docx")
            
            # 文本文件
            elif ext == ".txt":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            
            else:
                logger.warning(f"不支持的文件类型: {ext}，路径: {file_path}")
                raise Exception(f"不支持的文件类型: {ext}")
            
            return content.strip()
            
        except Exception as e:
            logger.error(f"读取文件 {file_path} 失败: {e}")
            raise

    def _mineru_api_process(self, file_path: str, file_type: str) -> str:
        """使用MinerU API处理文件"""
        try:
            import requests
        except ImportError:
            raise Exception("需要安装requests库来调用MinerU API: pip install requests")
            
        # 准备请求头
        headers = {
            "Authorization": f"Bearer {self.config.mineru_api_key}",
            "Accept": "application/json"
        }
        
        # 上传文件
        try:
            with open(file_path, "rb") as f:
                files = {"file": (os.path.basename(file_path), f, "application/octet-stream")}
                data = {
                    "is_ocr": True,  # 图片和扫描PDF需要OCR
                    "enable_formula": True,
                    "enable_table": True,
                    "language": "ch",
                    "model_version": "v2"
                }
                
                response = requests.post(
                    f"{self.config.mineru_api_url}/extract/task",
                    headers=headers,
                    files=files,
                    data=data,
                    timeout=30
                )
                
                response.raise_for_status()  # 检查HTTP错误状态
                task_data = response.json()
                
                if task_data.get("status") != "success":
                    raise Exception(f"MinerU API请求失败: {task_data.get('message', '未知错误')}")
                
                task_id = task_data["data"]["task_id"]
                logger.info(f"MinerU任务创建成功，任务ID: {task_id}")
                
        except requests.exceptions.RequestException as e:
            logger.error(f"MinerU API调用失败: {e}")
            raise Exception(f"MinerU API调用失败: {str(e)}")
        
        # 轮询获取结果
        max_attempts = 30  # 最多轮询30次
        retry_interval = 5  # 每5秒轮询一次
        
        for attempt in range(max_attempts):
            try:
                result_response = requests.get(
                    f"{self.config.mineru_api_url}/extract-results/{task_id}",
                    headers=headers,
                    timeout=10
                )
                
                result_response.raise_for_status()
                result_data = result_response.json()
                
                if result_data.get("status") == "completed":
                    return result_data["data"]["markdown_content"]
                elif result_data.get("status") == "failed":
                    error_msg = result_data.get("message", "未知错误")
                    logger.error(f"MinerU任务处理失败: {error_msg}")
                    raise Exception(f"MinerU任务处理失败: {error_msg}")
                
                # 任务仍在处理中
                logger.debug(f"MinerU任务 {task_id} 处理中（{attempt+1}/{max_attempts}）")
                time.sleep(retry_interval)
                
            except requests.exceptions.RequestException as e:
                logger.warning(f"获取MinerU任务结果失败（{attempt+1}/{max_attempts}）: {e}")
                time.sleep(retry_interval)
        
        # 任务超时
        raise Exception(f"MinerU任务 {task_id} 处理超时，请检查API状态")

    def _split_into_chunks(self, content: str, file_path: str) -> List[str]:
        """
        将文本内容分割成块
        
        Args:
            content: 要分割的文本内容
            file_path: 文件路径，用于生成块标题
            
        Returns:
            分块后的文本列表
        """
        if not content:
            return []
            
        chunks = []
        current_chunk = ""
        separators = self.config.separators
        
        # 提取文件名作为标题前缀
        file_name = os.path.basename(file_path)
        title_prefix = f"文件: {file_name}\n\n"
        
        # 按分隔符优先级进行分割
        for separator in separators:
            if len(content) <= self.config.chunk_size:
                break
                
            parts = content.split(separator)
            content = []
            for part in parts:
                if len(part) > self.config.chunk_size:
                    content.append(part)
                else:
                    if len(current_chunk) + len(part) + len(separator) <= self.config.chunk_size:
                        current_chunk += part + separator
                    else:
                        if current_chunk:
                            chunks.append(current_chunk)
                        current_chunk = part + separator
            if current_chunk:
                chunks.append(current_chunk)
                current_chunk = ""
            content = "\n".join(content)
        
        # 处理剩余内容
        if content:
            # 如果仍超过块大小，使用最后手段 - 按字符分割
            if len(content) > self.config.chunk_size:
                logger.warning(f"文件 {file_path} 内容无法按自然分隔符分割，将强制按字符分割")
                for i in range(0, len(content), self.config.chunk_size - self.config.chunk_overlap):
                    chunks.append(content[i:i+self.config.chunk_size])
            else:
                chunks.append(content)
        
        # 添加标题前缀和块编号
        numbered_chunks = []
        for i, chunk in enumerate(chunks, 1):
            chunk_header = f"{title_prefix}块 {i}/{len(chunks)}\n\n"
            numbered_chunk = chunk_header + chunk
            numbered_chunks.append(numbered_chunk)
            
        logger.info(f"文件 {file_path} 已分割为 {len(numbered_chunks)} 个块")
        return numbered_chunks

    def _get_embedding_with_retry(self, text: str) -> List[float]:
        """带重试机制的嵌入向量获取"""
        for attempt in range(self.config.retry_attempts):
            try:
                response = self.openai_client.embeddings.create(
                    input=text,
                    model=self.config.embedding_model
                )
                return response.data[0].embedding
                
            except Exception as e:
                if attempt < self.config.retry_attempts - 1:
                    logger.warning(f"嵌入API调用失败 (尝试 {attempt+1}/{self.config.retry_attempts})，{self.config.retry_delay}秒后重试: {e}")
                    time.sleep(self.config.retry_delay)
                else:
                    logger.error(f"嵌入API调用失败 (已达最大重试次数 {self.config.retry_attempts}): {e}")
                    raise

    def process_file(self, file_path: str) -> bool:
        """处理单个文件：读取内容、分块、生成向量、更新数据库"""
        try:
            # 确定文件所属主题
            topic = self._get_file_topic(file_path)
            
            # 获取文件元数据
            file_stat = os.stat(file_path)
            last_modified = datetime.fromtimestamp(file_stat.st_mtime)
            
            # 计算文件哈希
            file_hash = self._get_file_hash(file_path)
            
            # 检查数据库中是否已有相同哈希的文件
            current_status = self.db.get_file_status(file_path)
            if current_status and current_status["status"] == "success" and current_status.get("file_hash") == file_hash:
                logger.info(f"文件 {file_path} (主题: {topic}) 未发生变化，无需处理")
                return True
                
            # 更新状态为处理中（初始状态）
            self.db.upsert_file(
                file_path=file_path,
                topic=topic,
                file_hash=file_hash,
                last_modified=last_modified,
                status="processing"
            )
            
            # 读取文件内容
            content = self._read_file_content(file_path)
            if not content:
                logger.warning(f"文件内容为空，跳过向量化: {file_path}")
                # 更新状态为成功但内容为空
                self.db.upsert_file(
                    file_path=file_path,
                    topic=topic,
                    file_hash=file_hash,
                    last_modified=last_modified,
                    status="success",
                    total_chunks=0,
                    processed_chunks=0,
                    error_message="文件内容为空"
                )
                return True
                
            # 将内容分割成块
            chunks = self._split_into_chunks(content, file_path)
            total_chunks = len(chunks)
            
            if total_chunks == 0:
                logger.warning(f"文件 {file_path} 分块后无内容，跳过向量化")
                self.db.upsert_file(
                    file_path=file_path,
                    topic=topic,
                    file_hash=file_hash,
                    last_modified=last_modified,
                    status="success",
                    total_chunks=0,
                    processed_chunks=0
                )
                return True
                
            # 更新文件状态 - 开始分块处理
            file_id = self.db.upsert_file(
                file_path=file_path,
                topic=topic,
                file_hash=file_hash,
                last_modified=last_modified,
                status="processing",
                total_chunks=total_chunks,
                processed_chunks=0
            )
            
            # 处理每个块
            for i, chunk in enumerate(chunks, 1):
                try:
                    # 生成块嵌入向量
                    chunk_vector = self._get_embedding_with_retry(chunk)
                    
                    # 验证向量维度
                    if len(chunk_vector) != self.config.embedding_dim:
                        raise Exception(f"嵌入向量维度不匹配，模型返回 {len(chunk_vector)} 维，配置预期 {self.config.embedding_dim} 维")
                    
                    # 保存块向量到数据库
                    self.db.upsert_chunk(
                        file_id=file_id,
                        chunk_num=i,
                        total_chunks=total_chunks,
                        content=chunk,
                        vector=chunk_vector
                    )
                    
                    logger.info(f"文件块 {file_path}:{i}/{total_chunks} 处理成功")
                    
                except Exception as e:
                    logger.error(f"处理文件块 {file_path}:{i} 失败: {e}")
                    # 更新文件状态为失败
                    self.db.upsert_file(
                        file_path=file_path,
                        topic=topic,
                        file_hash=file_hash,
                        last_modified=last_modified,
                        status="failed",
                        total_chunks=total_chunks,
                        processed_chunks=i-1,  # 已处理的块数
                        error_message=f"块 {i} 处理失败: {str(e)}"
                    )
                    return False
            
            # 所有块处理完成，更新文件状态为成功
            self.db.upsert_file(
                file_path=file_path,
                topic=topic,
                file_hash=file_hash,
                last_modified=last_modified,
                status="success",
                total_chunks=total_chunks,
                processed_chunks=total_chunks
            )
            
            # 发送文件处理成功事件
            from ..utils.events import event_queue, get_event_loop
            event_data = {
                "type": "file_processed",
                "file_path": file_path,
                "topic": topic,
                "total_chunks": total_chunks,
                "timestamp": datetime.now().isoformat()
            }
            
            asyncio.run_coroutine_threadsafe(
                event_queue.put(event_data), 
                get_event_loop()
            )
            
            logger.info(f"文件 {file_path} (主题: {topic}) 分块处理完成，共 {total_chunks} 块")
            return True
            
        except Exception as e:
            error_msg = str(e)
            topic = self._get_file_topic(file_path) if 'topic' not in locals() else topic
            file_hash = self._get_file_hash(file_path) if 'file_hash' in locals() else ""
            
            try:
                # 获取文件修改时间
                if os.path.exists(file_path):
                    file_stat = os.stat(file_path)
                    last_modified = datetime.fromtimestamp(file_stat.st_mtime)
                else:
                    last_modified = datetime.now()
                    
                # 更新数据库，状态设为失败
                self.db.upsert_file(
                    file_path=file_path,
                    topic=topic,
                    file_hash=file_hash,
                    last_modified=last_modified,
                    status="failed",
                    error_message=error_msg
                )
                
            except Exception as db_error:
                logger.error(f"更新文件失败状态到数据库时出错: {db_error}")
            
            # 发送文件处理失败事件
            from ..utils.events import event_queue, get_event_loop
            event_data = {
                "type": "file_failed",
                "file_path": file_path,
                "topic": topic,
                "error": error_msg,
                "timestamp": datetime.now().isoformat()
            }
            
            asyncio.run_coroutine_threadsafe(
                event_queue.put(event_data), 
                get_event_loop()
            )
            
            logger.error(f"文件 {file_path} (主题: {topic}) 处理失败: {error_msg}")
            return False

    def process_batch(self, file_paths: List[str]) -> None:
        """批处理文件列表"""
        if not file_paths:
            return
            
        logger.info(f"开始批处理 {len(file_paths)} 个文件")
        for file_path in file_paths:
            try:
                self.process_file(file_path)
            except Exception as e:
                logger.error(f"处理文件 {file_path} 失败: {e}")
                continue
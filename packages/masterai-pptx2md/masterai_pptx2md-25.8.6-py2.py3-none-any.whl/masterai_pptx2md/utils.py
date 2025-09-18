import io
import base64
import logging
from pathlib import Path
from io import StringIO
from mimetypes import MimeTypes
from typing import Union, Optional
from PIL import Image
from typing import Optional, Tuple

import pillow_heif
import pillow_avif
pillow_heif.register_heif_opener()

from pptx.shapes.picture import Picture
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.text.text import Font
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR

from masterai_pptx2md.http_file_types import apache2_mime

logger = logging.getLogger(__name__)

my_mimetypes = MimeTypes()
my_mimetypes.readfp(StringIO(apache2_mime), True)


def is_title(shape: Union[SlidePlaceholder, Shape, GraphicFrame, Picture]) -> bool:
    if not shape.is_placeholder:
        return False

    if shape.placeholder_format.type in (
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.SUBTITLE,
        PP_PLACEHOLDER.VERTICAL_TITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
    ):
        return True
    else:
        return False


def is_text_block(shape: Union[SlidePlaceholder, Shape, GraphicFrame, Picture]) -> bool:
    if not shape.has_text_frame:
        return False

    if shape.is_placeholder and shape.placeholder_format.type in (PP_PLACEHOLDER.BODY,):
        return True

    if shape.text:  # type: ignore
        return True

    return False


def is_list_block(shape: Union[SlidePlaceholder, Shape, GraphicFrame, Picture]) -> bool:
    levels = []
    for para in shape.text_frame.paragraphs:  # type: ignore
        if para.level not in levels:
            levels.append(para.level)
        if para.level != 0 or len(levels) > 1:
            return True
    return False


def is_accent(font: Font) -> bool:
    if font.underline:
        return True
    if font.italic:
        return True
    if font.color.type == MSO_COLOR_TYPE.SCHEME and font.color.theme_color in (
        MSO_THEME_COLOR.ACCENT_1,
        MSO_THEME_COLOR.ACCENT_2,
        MSO_THEME_COLOR.ACCENT_3,
        MSO_THEME_COLOR.ACCENT_4,
        MSO_THEME_COLOR.ACCENT_5,
        MSO_THEME_COLOR.ACCENT_6,
    ):
        return True
    return False


def is_strong(font: Font) -> bool:
    if font.bold:
        return True
    if font.color.type == MSO_COLOR_TYPE.SCHEME and font.color.theme_color in (
        MSO_THEME_COLOR.DARK_1,
        MSO_THEME_COLOR.DARK_2,
    ):
        return True
    return False


def image_to_base64(image_bytes: bytes, ext: str) -> Optional[str]:
    mimetype, _ = my_mimetypes.guess_type(f"a.{ext}")
    if not mimetype:
        logger.error(f"ext: {ext} not find mimetype")
        return
    encoded_string: str = base64.b64encode(image_bytes).decode("utf-8")
    return f"data:{mimetype};base64,{encoded_string}"


def to_image(data: bytes) -> Optional[Image.Image]:
    try:
        return Image.open(io.BytesIO(data))
    except Exception as e:
        logger.warning(f"to_image fail {e}")


def image2png(content: bytes) -> Optional[bytes]:
    try:
        image: Image = Image.open(io.BytesIO(content))
        value = io.BytesIO()
        image.save(value, format="PNG")
        value.seek(0)
        return value.getvalue()
    except Exception as e:
        logger.warning(f"image2png fail {e}")
        return


def get_image_size(image_file_path: Optional[str] = None, image_data: Optional[bytes] = None) -> Tuple[Optional[int], Optional[int]]:
    if not image_file_path and not image_data:
        return None, None
    
    if image_file_path:
        image_file_path_data: bytes = open(image_file_path, 'rb').read()
        image_data = image_file_path_data

    if not image_data:
        return None, None

    image: Optional[Image] = to_image(image_data)
    if not image:
        return None, None

    return image.width, image.height
    

def check_image_size(min_image_width: Optional[int], min_image_height: Optional[int], width: Optional[int], height: Optional[int]) -> bool:
    if not width or not height:
        return False

    if min_image_width and min_image_width >= width:
        return False
    
    if min_image_height and min_image_height >= height:
        return False

    return True
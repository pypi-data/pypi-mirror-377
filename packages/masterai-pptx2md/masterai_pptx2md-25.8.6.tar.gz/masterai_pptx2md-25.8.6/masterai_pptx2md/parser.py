import io
import os
import oss2
import logging
from tqdm import tqdm
from hashlib import md5
from datetime import date
from operator import attrgetter
from collections import defaultdict
from typing import List, Union, Optional, IO, Dict, Set

from pptx.slide import Slide
from pptx.shapes.picture import Picture
from pptx.shapes.autoshape import Shape
from pptx.presentation import Presentation as PresentationClass
from pptx.api import Presentation as PresentationFunc
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.shapetree import SlideShapes
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.text.text import _Paragraph as Paragraph, _Run as Run

from masterai_pptx2md.models import Config
from masterai_pptx2md.outputter import MarkDownOutPutter
from masterai_pptx2md.utils import (
    is_title,
    is_text_block,
    is_list_block,
    is_accent,
    is_strong,
    image_to_base64,
)
from masterai_pptx2md.exceptions import ConfigException
from masterai_pptx2md.utils import image2png, check_image_size, get_image_size


logger = logging.getLogger(__name__)


class Parse:

    def __init__(self, config: Config, out_putter: MarkDownOutPutter) -> None:
        self.config = config
        self.out_putter = out_putter
        self.oss_client: Optional[oss2.Bucket]
        self.slide_pics_map: Dict[int, Set] = defaultdict(set)  # 用于图片去重
        self.oss_client = None
        self.storage = None
        if not self.config.disable_image and self.config.upload_image:
            if not self.config.oss_config and not self.config.storage:
                raise ConfigException(f"miss oss and storage config while upload_image is true")
            if self.config.oss_config:
                if os.environ.get("LOCAL_RUN"):
                    endpoint: str = self.config.oss_config.endpoint_public
                else:
                    endpoint: str = self.config.oss_config.endpoint
                self.oss_client = oss2.Bucket(
                    auth=oss2.Auth(
                        access_key_id=self.config.oss_config.access_key_id,
                        access_key_secret=self.config.oss_config.access_key_secret,
                    ),
                    endpoint=endpoint,
                    bucket_name=self.config.oss_config.bucket_name,
                )
            else:
                self.storage = self.config.storage

    def upload_image(self, image_id: str, image_bytes: bytes, suffix: str) -> Optional[str]:
        filename: str = f"{image_id}.{suffix}"

        if self.config.oss_config and self.oss_client:
            key: str = os.path.join(
                f"{self.config.oss_config.prefix.strip('/')}/{date.today().strftime('%Y%m%d')}",
                filename,
            )
            url: str = os.path.join(self.config.oss_config.cdn_host, key)
            if self.oss_client.object_exists(key):
                return url
            for _ in range(3):
                try:
                    self.oss_client.put_object(key, image_bytes)
                    return url
                except Exception as e:
                    continue
        elif self.storage:
            url: str = self.storage.upload_by_name(filename, image_bytes)
            return url
        else:
            raise ConfigException(f"oss client is none")
        return None

    def get_formatted_text(self, para: Paragraph) -> str:
        res: str = ""
        for run in para.runs:
            run: Run
            text: str = run.text
            if text == "":
                continue
            if not self.config.disable_escaping:
                text: str = self.out_putter.get_escaped(text)
            try:
                if run.hyperlink.address:
                    text: str = self.out_putter.get_hyperlink(
                        text, run.hyperlink.address
                    )
            except:
                text = self.out_putter.get_hyperlink(
                    text, "error:ppt-link-parsing-issue"
                )
            if is_accent(run.font):
                text: str = self.out_putter.get_accent(text)
            elif is_strong(run.font):
                text: str = self.out_putter.get_strong(text)
            if not self.config.disable_color:
                if run.font.color.type == MSO_COLOR_TYPE.RGB:
                    text: str = self.out_putter.get_colored(text, run.font.color.rgb)
            res += text
        return res.strip()

    def process_title(
        self, shape: Union[SlidePlaceholder, Shape], slide_idx: int
    ) -> None:
        if shape.has_text_frame:
            text: str = shape.text_frame.text.strip()
            self.out_putter.put_title(text, 1)

    def process_text_block(self, shape: Union[SlidePlaceholder, Shape], _: int) -> None:
        if is_list_block(shape):
            # generate list block
            for para in shape.text_frame.paragraphs:
                para: Paragraph
                if para.text.strip() == "":
                    continue
                text: str = self.get_formatted_text(para)
                self.out_putter.put_list(text, para.level)
            self.out_putter.write("\n")
        else:
            # generate paragraph block
            for para in shape.text_frame.paragraphs:
                para: Paragraph
                if para.text.strip() == "":
                    continue
                text: str = self.get_formatted_text(para)
                self.out_putter.put_para(text)

    def process_notes(self, text: str, _: int) -> None:
        self.out_putter.put_para("---")
        self.out_putter.put_para(text)

    def process_picture(self, shape: Picture, slide_idx: int) -> None:
        if self.config.disable_image:
            return

        pic_ext: str = shape.image.ext  # bmp gif jpg png tiff wmf
        if self.config.allow_image_format and pic_ext not in self.config.allow_image_format:
            return

        image_bytes: bytes = shape.image.blob
        if pic_ext in {"bmp", "heic", "heif", "tiff"}:
            image_bytes_png: Optional[bytes] = image2png(image_bytes)
            if not image_bytes_png:
                return
            image_bytes = image_bytes_png
            pic_ext = "png"
        
        width, height = get_image_size(image_data=image_bytes)
        if not check_image_size(self.config.min_image_width, self.config.min_image_height, width=width, height=height):
            return

        image_id: str = md5(image_bytes).hexdigest()
        if self.config.skip_duplicate_image:
            if image_id not in self.slide_pics_map[slide_idx]:
                self.slide_pics_map[slide_idx].add(image_id)
            else:
                return
        image_content: Optional[str]
        if self.config.upload_image:
            image_content = self.upload_image(image_id, image_bytes, pic_ext)
        else:
            image_content = image_to_base64(image_bytes, pic_ext)
        if not image_content:
            return
        self.out_putter.put_image(image_content, self.config.max_img_width, width, height, suffix=pic_ext)

    def process_table(self, shape: GraphicFrame, _: int) -> None:
        table: List[List[str]] = [
            [cell.text for cell in row.cells] for row in shape.table.rows
        ]
        if len(table) > 0:
            self.out_putter.put_table(table)

    def ungroup_shapes(
        self, shapes: SlideShapes
    ) -> List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]]:
        res: List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]] = []
        for shape in shapes:  # type: ignore
            shape: Union[SlidePlaceholder, Shape, GraphicFrame, Picture]
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    res.extend(self.ungroup_shapes(shape.shapes))  # type: ignore
                else:
                    res.append(shape)
            except Exception as e:
                logger.error(f"failed to load shape {shape}, skipped. error: {e}")
        return res

    def process_shapes(
        self,
        current_shapes: List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]],
        slide_id: int,
    ) -> None:
        for shape in current_shapes:
            shape: Union[SlidePlaceholder, Shape, GraphicFrame, Picture]
            if is_title(shape):
                try:
                    self.process_title(shape, slide_id + 1)  # type: ignore
                except Exception as e:
                    logger.exception(
                        f"failed to process_title {shape}, skipped. error: {e}"
                    )
            elif is_text_block(shape):
                try:
                    self.process_text_block(shape, slide_id + 1)  # type: ignore
                except Exception as e:
                    logger.exception(
                        f"failed to process_text_block {shape}, skipped. error: {e}"
                    )
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    self.process_picture(shape, slide_id + 1)  # type: ignore
                except Exception as e:
                    logger.error(f"Failed to process picture, skipped: {e}")
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                self.process_table(shape, slide_id + 1)  # type: ignore
            else:
                try:
                    ph = shape.placeholder_format
                    if (
                        ph.type == PP_PLACEHOLDER.OBJECT
                        and hasattr(shape, "image")
                        and getattr(shape, "image")
                    ):
                        self.process_picture(shape, slide_id + 1)  # type: ignore
                except:
                    pass

    def parse(self, pptx_content: bytes) -> str:
        pptx: IO[bytes] = io.BytesIO(pptx_content)
        prs: PresentationClass = PresentationFunc(pptx)
        for idx, slide in enumerate(tqdm(prs.slides, desc="Converting slides")):
            idx: int
            slide: Slide
            shapes: List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]] = []
            try:
                shapes_with_none: List[
                    Union[SlidePlaceholder, Shape, GraphicFrame, Picture]
                ] = self.ungroup_shapes(slide.shapes)
                shapes: List[Union[SlidePlaceholder, Shape, GraphicFrame, Picture]] = (
                    sorted(
                        shapes_with_none,
                        key=lambda x: (
                            getattr(x, "top", 0) or 0,
                            getattr(x, "left", 0) or 0,
                        ),
                    )
                )
            except:
                logger.error(
                    "Bad shapes encountered in this slide. Please check or move them and try again."
                )
                logger.error("shapes:")
                try:
                    for sp in slide.shapes:
                        logger.error(sp.shape_type)
                        logger.error(sp.top, sp.left, sp.width, sp.height)
                except:
                    logger.error("failed to print all bad shapes.")

            self.process_shapes(shapes, idx + 1)

            if not self.config.disable_notes and slide.has_notes_slide:
                if not slide.notes_slide or not slide.notes_slide.notes_text_frame:
                    continue
                text: str = slide.notes_slide.notes_text_frame.text  # type: ignore
                if text:
                    self.process_notes(text, idx + 1)
            if idx < len(prs.slides) - 1 and self.config.enable_slides:
                self.out_putter.put_para("\n---\n")

        return self.out_putter.read()
